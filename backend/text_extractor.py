from pathlib import Path

def extract_text(file_path: str, file_ext: str) -> str:
    ext = file_ext.lower()
    if ext == ".txt": return _extract_txt(file_path)
    elif ext == ".pdf": return _extract_pdf(file_path)
    elif ext == ".docx": return _extract_docx(file_path)
    elif ext in (".pptx", ".ppt"): return _extract_pptx(file_path)
    else: raise ValueError(f"Unsupported format: {ext}")

def _extract_txt(file_path):
    for enc in ["utf-8", "latin-1", "cp1252"]:
        try:
            with open(file_path, "r", encoding=enc) as f: return f.read()
        except UnicodeDecodeError: continue
    raise ValueError("Could not decode text file.")

def _extract_pdf(file_path):
    text_parts = []
    try:
        import PyPDF2
        with open(file_path, "rb") as f:
            reader = PyPDF2.PdfReader(f)
            for page in reader.pages:
                t = page.extract_text()
                if t: text_parts.append(t)
    except Exception: pass
    if not text_parts:
        try:
            import pdfplumber
            with pdfplumber.open(file_path) as pdf:
                for page in pdf.pages:
                    t = page.extract_text()
                    if t: text_parts.append(t)
        except Exception: pass
    return "\n\n".join(text_parts)

def _extract_docx(file_path):
    try:
        from docx import Document
        doc = Document(file_path)
        paragraphs = [p.text for p in doc.paragraphs if p.text.strip()]
        for table in doc.tables:
            for row in table.rows:
                for cell in row.cells:
                    if cell.text.strip(): paragraphs.append(cell.text.strip())
        return "\n\n".join(paragraphs)
    except ImportError: raise ValueError("Run: pip install python-docx")
    except Exception as e: raise ValueError(f"DOCX error: {e}")

def _extract_pptx(file_path):
    try:
        from pptx import Presentation
        prs = Presentation(file_path)
        slides_text = []
        for i, slide in enumerate(prs.slides, 1):
            content = [f"--- Slide {i} ---"]
            for shape in slide.shapes:
                if hasattr(shape, "text") and shape.text.strip():
                    content.append(shape.text.strip())
            if len(content) > 1: slides_text.append("\n".join(content))
        return "\n\n".join(slides_text)
    except ImportError: raise ValueError("Run: pip install python-pptx")
    except Exception as e: raise ValueError(f"PPTX error: {e}")
